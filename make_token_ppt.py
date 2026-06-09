from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── 색상 ─────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x17, 0x2A)
RED     = RGBColor(0xFF, 0x44, 0x44)
RED2    = RGBColor(0xCC, 0x22, 0x22)
ORANGE  = RGBColor(0xFF, 0x99, 0x22)
GREEN   = RGBColor(0x22, 0xCC, 0x66)
GREEN2  = RGBColor(0x16, 0x8A, 0x45)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x88, 0x99, 0xAA)
YELLOW  = RGBColor(0xFF, 0xDD, 0x00)
DARK    = RGBColor(0x16, 0x24, 0x3E)
DARKER  = RGBColor(0x0A, 0x10, 0x1E)

def rect(x, y, w, h, fill, line=None, lw=1):
    sh = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def txt(x, y, w, h, text, size=9, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def bar(x, y, w, h, fill, pct_label, tok_label, color=WHITE):
    rect(x, y, w, h, fill)
    txt(x + w + 0.3, y + h/2 - 0.35, 6, 0.7, tok_label, size=10, bold=True, color=color)

# ── 배경 ─────────────────────────────────────────────
rect(0, 0, 33.87, 19.05, BG)

# ── 제목 ─────────────────────────────────────────────
rect(0, 0, 33.87, 1.3, RGBColor(0x08, 0x10, 0x20))
txt(0.5, 0.2, 32, 0.9,
    "AI 토큰 사용량 비교  —  최적화 전 vs 현재 적용",
    size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  섹션 1: 법령 데이터 처리 (왼쪽)
# ══════════════════════════════════════════════════════════════
rect(0.4, 1.5, 15.8, 10.5, DARK, line=RGBColor(0x44,0x66,0x99))
txt(0.7, 1.65, 14, 0.6, "① 법령 데이터 처리", size=11, bold=True, color=YELLOW)

# 라벨
txt(0.7, 2.5, 6.5, 0.5, "그냥 썼을 때", size=9, bold=True, color=RED)
txt(0.7, 5.6, 6.5, 0.5, "지금 적용 (MCP 파싱)", size=9, bold=True, color=GREEN)

# Before 막대 - 500,000 토큰 (법령 XML 직접)
rect(0.7, 3.1, 13.8, 1.8, RED2)
txt(0.9, 3.3, 13, 0.55, "law.go.kr XML 원본  →  AI에 직접 전달", size=9, color=WHITE)
txt(0.9, 3.85, 13, 0.7, "500,000 토큰", size=14, bold=True, color=WHITE)
txt(0.7, 5.0, 13.8, 0.5, "⚠️  Claude 200k 컨텍스트 초과  →  처리 자체가 불가능", size=8.5, color=ORANGE)

# After 막대 - 600 토큰
rect(0.7, 6.2, 0.17, 1.8, GREEN2)   # 좁은 막대 (600/500000 비율)
rect(0.7, 6.2, 13.8, 1.8, RGBColor(0x0A,0x30,0x18), line=GREEN, lw=1)
txt(0.9, 6.4, 10, 0.55, "XML 2.3MB  →  파싱  →  1.8KB JSON만 전달", size=9, color=WHITE)
txt(0.9, 6.95, 10, 0.7, "600 토큰", size=14, bold=True, color=GREEN)

# 절감 뱃지
rect(9.5, 6.5, 4.3, 1.3, RGBColor(0x0A, 0x40, 0x20), line=GREEN)
txt(9.7, 6.65, 3.9, 0.9, "99.9% 절감\n833배 차이", size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 화살표 설명
txt(0.7, 8.3, 15, 0.5, "2,341,547 bytes  →  1,833 bytes  /  41,737줄  →  36줄  /  1일 캐시 적용", size=8, color=GRAY)

# ══════════════════════════════════════════════════════════════
#  섹션 2: AI 챗봇 전체 입력 (오른쪽)
# ══════════════════════════════════════════════════════════════
rect(16.7, 1.5, 16.8, 10.5, DARK, line=RGBColor(0x44,0x66,0x99))
txt(17.0, 1.65, 15, 0.6, "② lawChat  —  요청당 AI 입력 토큰", size=11, bold=True, color=YELLOW)

# Before 레이블
txt(17.0, 2.5, 7, 0.5, "그냥 썼을 때", size=9, bold=True, color=RED)

# Before 스택 막대 (24,300 토큰)
BAR_W = 15.0
SCALE = BAR_W / 25000   # 25000이 최대

# 24300 총합
w_kb = 23000 * SCALE   # Obsidian 23,000
w_law = 300 * SCALE    # 법령 300
w_sys = 499 * SCALE    # 시스템 499
w_his = 500 * SCALE    # 히스토리 500

bx = 17.0
rect(bx,           3.1, w_kb,            1.6, RED2)
rect(bx + w_kb,    3.1, w_law,           1.6, RGBColor(0xAA,0x44,0x00))
rect(bx + w_kb + w_law, 3.1, w_sys,      1.6, RGBColor(0x66,0x44,0x00))
rect(bx + w_kb + w_law + w_sys, 3.1, w_his, 1.6, RGBColor(0x44,0x33,0x00))

txt(17.0, 4.8, 15, 0.45, "Obsidian 35노드 23,000  +  법령 300  +  시스템 499  +  히스토리 500", size=7.5, color=GRAY)
txt(17.0, 5.3, 7, 0.7, "총  24,300 토큰", size=13, bold=True, color=RED)

# After 레이블
txt(17.0, 6.3, 7, 0.5, "지금 적용 (RAG 최적화)", size=9, bold=True, color=GREEN)

# After 스택 막대 (2,800 토큰)
w_kb2  = 1500 * SCALE
w_law2 = 300  * SCALE
w_sys2 = 499  * SCALE
w_his2 = 500  * SCALE

rect(bx,                     7.0, w_kb2,  1.6, GREEN2)
rect(bx + w_kb2,             7.0, w_law2, 1.6, RGBColor(0x00,0x66,0x44))
rect(bx + w_kb2 + w_law2,    7.0, w_sys2, 1.6, RGBColor(0x00,0x44,0x44))
rect(bx + w_kb2+w_law2+w_sys2, 7.0, w_his2, 1.6, RGBColor(0x00,0x33,0x44))

txt(17.0, 8.65, 15, 0.45, "Obsidian 5노드 1,500  +  법령 300  +  시스템 499  +  히스토리 500", size=7.5, color=GRAY)
txt(17.0, 9.2, 7, 0.7, "총  2,800 토큰", size=13, bold=True, color=GREEN)

# 절감 뱃지
rect(26.5, 7.2, 5.5, 1.4, RGBColor(0x0A, 0x40, 0x20), line=GREEN)
txt(26.7, 7.38, 5.1, 1.0, "88% 절감\n8.7배 차이", size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  섹션 3: 범례
# ══════════════════════════════════════════════════════════════
rect(0.4, 12.3, 32.8, 0.9, RGBColor(0x10,0x1A,0x30))
txt(0.7, 12.45, 32, 0.6,
    "■ Obsidian 지식     ■ law.go.kr 조문     ■ 시스템 프롬프트     ■ 대화 히스토리",
    size=8.5, color=GRAY)

# ══════════════════════════════════════════════════════════════
#  섹션 4: 전체 요약 카드 3개
# ══════════════════════════════════════════════════════════════
txt(0.4, 13.4, 32, 0.6, "전체 요약", size=11, bold=True, color=YELLOW)

# 카드 1
rect(0.4, 14.1, 10.2, 4.4, RGBColor(0x20,0x08,0x08), line=RED)
txt(0.7, 14.3, 9.6, 0.6, "최적화 없이 쓴다면", size=10, bold=True, color=RED)
txt(0.7, 15.05, 9.6, 3.0,
    "법령 XML 전체\n→ 500,000 토큰\n→ 컨텍스트 초과\n→ 처리 불가  ❌",
    size=10, color=WHITE)

# 카드 2
rect(11.3, 14.1, 10.5, 4.4, RGBColor(0x20,0x14,0x00), line=ORANGE)
txt(11.6, 14.3, 9.9, 0.6, "RAG만 적용 (최적화 전)", size=10, bold=True, color=ORANGE)
txt(11.6, 15.05, 9.9, 3.0,
    "Obsidian 35노드 폭발\n→ 24,300 토큰\n→ 동작하지만 낭비\n→ 개선 필요  ⚠️",
    size=10, color=WHITE)

# 카드 3
rect(22.6, 14.1, 10.8, 4.4, RGBColor(0x04,0x20,0x10), line=GREEN)
txt(22.9, 14.3, 10.2, 0.6, "현재 적용 (MCP + RAG 최적화)", size=10, bold=True, color=GREEN)
txt(22.9, 15.05, 10.2, 3.0,
    "파싱 + 노드 제한\n→ 2,800 토큰\n→ 원본 대비 99% 절감\n→ 최적 상태  ✅",
    size=10, color=WHITE)

# ── 하단 강조 ─────────────────────────────────────────────
rect(0, 18.55, 33.87, 0.5, RGBColor(0x08, 0x10, 0x20))
txt(0.5, 18.6, 32.5, 0.38,
    "법령 원본 500,000 토큰  →  MCP 파싱 600 토큰  (99.9% 절감)   |   "
    "lawChat 24,300 토큰  →  2,800 토큰  (88% 절감)   |   종합 절감률 약 99%",
    size=8, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

out = "/Users/kimjongjin/coding/calculator-main/KTL_토큰절감_비교.pptx"
prs.save(out)
print("저장:", out)
