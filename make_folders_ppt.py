# -*- coding: utf-8 -*-
"""KTL Workspace 폴더별 역할 — 1장 쉬운 설명 PPT."""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
SW, SH = 33.87, 19.05

C_BG     = RGBColor(0x0B, 0x12, 0x22)
C_DARK   = RGBColor(0x10, 0x1B, 0x30)
C_PARSER = RGBColor(0x2F, 0x6F, 0xD6)
C_CALC   = RGBColor(0x1F, 0x9D, 0x8A)
C_LAW    = RGBColor(0x9B, 0x59, 0xB6)
C_INFRA  = RGBColor(0xC2, 0x96, 0x2C)
C_ANTI   = RGBColor(0x6D, 0x5A, 0xE0)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_INK    = RGBColor(0x0B, 0x12, 0x22)
C_SUB    = RGBColor(0xAA, 0xBC, 0xD0)
C_LINE   = RGBColor(0x2C, 0x3C, 0x5A)

s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background(); bg.shadow.inherit = False

def box(x, y, w, h, fill, line=None, lw=1.0, radius=True):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shp, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(x, y, w, h, text, size, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(0.08); tf.margin_top = tf.margin_bottom = Cm(0.02)
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = '맑은 고딕'
    return tb

def card(x, y, w, h, accent, name, why):
    box(x, y, w, h, C_DARK, line=C_LINE)
    box(x, y, 0.15, h, accent, radius=False)
    txt(x+0.35, y+0.15, w-0.55, 0.6, name, 11, C_WHITE, True)
    txt(x+0.35, y+0.78, w-0.55, h-0.85, "└ " + why, 9, C_SUB, spacing=1.05)

def group(x, y, w, h, accent, head, cards):
    box(x, y, w, h, C_BG, line=accent, lw=1.5)
    box(x, y, w, 0.95, accent, radius=True)
    txt(x, y+0.05, w, 0.9, head, 13, C_INK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    cy = y+1.2; gap = 0.25; ch = (h-1.35-gap*(len(cards)-1))/len(cards)
    for nm, why in cards:
        card(x+0.35, cy, w-0.7, ch, accent, nm, why)
        cy += ch + gap

# ── 타이틀 ──
box(0, 0, 0.5, SH, C_ANTI, radius=False)
txt(1.0, 0.5, 31, 1.0, "KTL Workspace 프로젝트 한눈에 — 폴더별 역할", 23, C_WHITE, True)
txt(1.0, 1.6, 31, 0.7, "사진 수집 ▶ 정밀 계산 ▶ 법령 해석  +  이를 떠받치는 인프라·도구", 12, C_SUB)

# ── 파이프라인 3그룹 (상단) ──
GY, GH = 2.7, 8.4
group(0.9, GY, 13.2, GH, C_PARSER, "① 데이터 수집 · OCR", [
    ("📱 photo-ocr-app", "사진 한 장으로 측정값 자동 추출 (사람이 안 침)"),
    ("🗄 parser-photo-server", "사진·OCR 데이터 창고 + 교환소 (오프라인 저장·조회)"),
    ("🤖 parser-orchestrator", "AI(Gemini) 호출을 순서대로 안전하게 (한도초과·실패 방지)"),
])
group(14.4, GY, 8.5, GH, C_CALC, "② 정밀 계산 · 판정", [
    ("🧮 calculator-main", "추출된 값으로 정도검사 판정\n(반복성·드리프트·직선성·현장적용계수)\n엑셀 기준 그대로 + AI 법령 챗봇 임베디드"),
])
group(23.2, GY, 9.7, GH, C_LAW, "③ 법령 해석 챗봇", [
    ("⚖️ law-api-vercel", "환경부 고시·물환경보전법을 검색용 지식으로 가공해 챗봇에 공급"),
    ("🔌 law-mcp-server", "AI가 법령을 직접 조회하게 노출 (MCP)"),
])

# ── 인프라·도구 (하단 6개) ──
IY = 11.9
box(0.9, IY, 32.0, 6.3, C_BG, line=C_INFRA, lw=1.5)
box(0.9, IY, 32.0, 0.9, C_INFRA, radius=True)
txt(0.9, IY+0.0, 32.0, 0.9, "④ 인프라 · 관리 · 도구 — 파이프라인을 떠받치는 받침대", 13, C_INK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
tools = [
    (C_INFRA, "💻 mac-studio-server", "로컬 동기화·SQLite 백업·Vercel 프록시 중계"),
    (C_INFRA, "📅 schedule-manager", "사업장 정도검사 일정 관리 대시보드"),
    (C_ANTI,  "🛡 harness", "AI 워커(Claude·Codex) 승인·격리 통제 샌드박스"),
    (C_CALC,  "📤 subunsuk-sender", "수분석 결과 Claydox 전송 + 카카오 알림"),
    (C_INFRA, "🔒 수분석 내부망 exe", "인터넷 차단 내부망용 오프라인 앱 (Electron)"),
    (C_INFRA, "🗓 yearlyplanner-ai", "연간 계획 관리 도구"),
]
cols, gx, gy0 = 3, 1.25, IY+1.15
cw = (32.0-0.7-(cols-1)*0.5)/cols; chh = 2.35
for i, (ac, nm, why) in enumerate(tools):
    r, c = divmod(i, cols)
    card(gx + c*(cw+0.5), gy0 + r*(chh+0.35), cw, chh, ac, nm, why)

txt(1.0, 18.35, 31, 0.6, "흐름 한 줄  ·  사진 → parser-photo-server 적재 → orchestrator가 AI로 추출 → calculator가 판정 → 법령 챗봇이 근거 설명",
    10, RGBColor(0x4A,0xDE,0x80), True)

out = "KTL_폴더별_역할.pptx"
prs.save(out)
print("저장:", out, "| 슬라이드", len(prs.slides._sldIdLst))
