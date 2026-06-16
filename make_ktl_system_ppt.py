# -*- coding: utf-8 -*-
"""KTL Workspace 통합 자동화 시스템 — 1장 설계 청사진 PPT."""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
SW, SH = 33.87, 19.05

C_BG     = RGBColor(0x0B, 0x12, 0x22)
C_PANEL  = RGBColor(0x14, 0x1F, 0x36)
C_DARK   = RGBColor(0x0E, 0x17, 0x2A)
C_PARSER = RGBColor(0x2F, 0x6F, 0xD6)
C_CALC   = RGBColor(0x1F, 0x9D, 0x8A)
C_LAW    = RGBColor(0x9B, 0x59, 0xB6)
C_ANTI   = RGBColor(0x6D, 0x5A, 0xE0)
C_CLAUDE = RGBColor(0xD9, 0x7B, 0x3A)
C_CODEX  = RGBColor(0x1F, 0x9D, 0x8A)
C_HUMAN  = RGBColor(0xC2, 0x96, 0x2C)
C_INFRA  = RGBColor(0x33, 0x44, 0x60)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_INK    = RGBColor(0x0B, 0x12, 0x22)
C_SUB    = RGBColor(0x9F, 0xB3, 0xC8)
C_LINE   = RGBColor(0x2C, 0x3C, 0x5A)
C_ACCENT = RGBColor(0x4A, 0xDE, 0x80)

s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = C_BG; bg.line.fill.background(); bg.shadow.inherit = False

def box(x, y, w, h, fill, line=None, radius=True, lw=1.0):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shp, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(x, y, w, h, text, size, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(0.08); tf.margin_top = tf.margin_bottom = Cm(0.02)
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = '맑은 고딕'
    return tb

def conn(x1, y1, x2, y2, c, w=2.5):
    ln = s.shapes.add_connector(2, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    ln.line.color.rgb = c; ln.line.width = Pt(w)

def arrow(cx, cy, c):  # ▶ 오른쪽 화살표
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(cx), Cm(cy), Cm(0.75), Cm(0.7))
    a.fill.solid(); a.fill.fore_color.rgb = c; a.line.fill.background(); a.shadow.inherit = False

def subbox(x, y, w, h, accent, name, desc):
    box(x, y, w, h, C_DARK, line=C_LINE)
    box(x, y, 0.13, h, accent, radius=False)
    txt(x+0.3, y+0.12, w-0.45, 0.55, name, 9.5, C_WHITE, True)
    txt(x+0.3, y+0.7, w-0.45, h-0.8, desc, 7.8, C_SUB, spacing=1.0)

# ── 타이틀 ──
box(0, 0, 0.5, SH, C_ANTI, radius=False)
txt(1.0, 0.45, 27, 1.0, "KTL Workspace 통합 자동화 시스템 — 설계 청사진", 22, C_WHITE, True)
txt(1.0, 1.55, 32, 0.7, "현장 OCR 수집(Parser)  ▶  정밀 계산·판정(Calculator)  ▶  법령 지식 서비스(Law API)  |  멀티에이전트 하네스 통제", 11.5, C_SUB)

# ════ 파이프라인 3단계 ════
PY, PH = 2.7, 6.7
def stage(x, w, c, head, items):
    box(x, PY, w, PH, C_PANEL, line=c, lw=1.5)
    txt(x, PY+0.18, w, 0.7, head, 11, c, True, PP_ALIGN.CENTER)
    iy = PY+1.05; gap = 0.22; ih = (PH-1.25-gap*(len(items)-1))/len(items)
    for nm, ds in items:
        subbox(x+0.35, iy, w-0.7, ih, c, nm, ds)
        iy += ih + gap

stage(0.9, 10.3, C_PARSER, "①  데이터 수집 · OCR 추출  [React/TS+Python]", [
    ("photo-ocr-app", "React 대시보드 — 사진·성적서 OCR 추출·분석 (Kakao/구조검사/먹는물/그래프)"),
    ("parser-orchestrator", "Python LLM 파이프라인 — Gemini 호출순서·전처리·레이트리밋 제어"),
    ("parser-photo-server", "Node/Express/SQLite — 사진·OCR 원본 적재 DB (PM2)"),
])
stage(11.6, 10.3, C_CALC, "②  정밀 계산 · 검사  [Vite+Serverless]", [
    ("calculator-main", "TMS·먹는물 정도검사 판정 엔진 + AI 챗봇 임베디드 (precision.js 순수계산)"),
    ("Version11_2026.xlsx", "계산 기준 DB — 수식·기준값 SSOT (sync-excel.py 자동 동기화)"),
])
stage(22.3, 10.7, C_LAW, "③  법령 해석 챗봇  [Node+Vercel]", [
    ("law-api-vercel", "법령 그래프 빌더 — 법제처 API 스파이더 + 조항 포함·인용 search-index"),
    ("AI 법령 챗봇", "WebRTC 음성인식(STT) — calculator 정도검사 UI 임베디드"),
])
# 입력 + 화살표
txt(0.9, 2.05, 10.3, 0.6, "📷 현장·카카오톡 (검사 성적서·현장 사진) ▼", 9.5, C_HUMAN, True, PP_ALIGN.CENTER)
arrow(11.0, 5.7, C_ACCENT)   # ①→②
arrow(21.7, 5.7, C_LAW)      # ③↔②

# ════ 인프라 밴드 ════
IY = 9.7
box(0.9, IY, 32.0, 1.85, C_INFRA, line=C_LINE)
txt(1.15, IY+0.12, 10, 0.55, "④  로컬 인프라 · 관리", 10.5, C_WHITE, True)
inf = [("mac-studio-server", "파일동기화·SQLite 백업·Vercel 프록시 중계"),
       ("schedule-manager", "사업장 정도검사 일정 관리 대시보드 (React/TS)"),
       ("수분석 내부망 exe", "인터넷 차단 내부망 오프라인 배포본 (Electron)")]
ix = 1.15
for nm, ds in inf:
    box(ix, IY+0.72, 10.35, 1.0, C_DARK, line=C_LINE)
    txt(ix+0.28, IY+0.78, 10, 0.5, nm, 9.3, C_ACCENT, True)
    txt(ix+0.28, IY+1.22, 10, 0.45, ds, 7.8, C_SUB)
    ix += 10.6

# ════ 하네스 통제 밴드 ════
HY = 11.95
box(0.9, HY, 32.0, 3.0, C_PANEL, line=C_ANTI, lw=1.5)
txt(1.15, HY+0.12, 31, 0.55, "🧭  멀티에이전트 하네스(Harness) — 인간 승인 게이트 하에 오케스트레이터가 기획, 워커가 코드 작성", 11, C_ANTI, True)
# 인간
box(1.3, HY+0.85, 6.6, 1.9, C_HUMAN); txt(1.3, HY+1.0, 6.6, 0.7, "👤 인간 개발자", 11, C_INK, True, PP_ALIGN.CENTER); txt(1.3, HY+1.7, 6.6, 1.0, "승인 게이트 — 명시 승인 없이\n외부 AI 자동 실행 불가", 8.3, C_INK, False, PP_ALIGN.CENTER)
# 오케스트레이터
box(8.6, HY+0.85, 8.2, 1.9, C_ANTI); txt(8.6, HY+1.0, 8.2, 0.7, "🧭 오케스트레이터 · Antigravity", 11, C_WHITE, True, PP_ALIGN.CENTER); txt(8.6, HY+1.7, 8.2, 1.0, "Gemini 3.1 Pro — 기획·마일스톤\ncall_worker.sh로 토큰 최소 노출", 8.3, C_WHITE, False, PP_ALIGN.CENTER)
# 워커
box(17.5, HY+0.85, 7.4, 1.9, C_CLAUDE); txt(17.5, HY+1.0, 7.4, 0.7, "🟧 워커 · claude-main", 10.5, C_INK, True, PP_ALIGN.CENTER); txt(17.5, HY+1.7, 7.4, 1.0, "코드 변경·작성 (추론·구현)\ntasks/ 내에서만 쓰기", 8.3, C_INK, False, PP_ALIGN.CENTER)
box(25.5, HY+0.85, 7.2, 1.9, C_CODEX); txt(25.5, HY+1.0, 7.2, 0.7, "🟦 워커 · codex-main", 10.5, C_WHITE, True, PP_ALIGN.CENTER); txt(25.5, HY+1.7, 7.2, 1.0, "코드 변경·작성 (코딩)\nAppend-Only 감사 로그", 8.3, C_WHITE, False, PP_ALIGN.CENTER)
arrow(7.95, HY+1.5, C_HUMAN); arrow(16.85, HY+1.5, C_ANTI)

# ════ 차별성 푸터 ════
FY = 15.25
def feat(x, w, c, t, b):
    box(x, FY, w, 3.25, C_DARK, line=c, lw=1.3); box(x, FY, w, 0.13, c, radius=False)
    txt(x+0.4, FY+0.3, w-0.7, 0.7, t, 12.5, c, True)
    txt(x+0.4, FY+1.15, w-0.7, 2.0, b, 9.8, C_WHITE, spacing=1.12)
feat(0.9, 10.4, C_PARSER, "🔌 온/오프라인 하이브리드",
     "온라인=Vercel 클라우드(AI 법령·서버 동기화), 오프라인=내부망 EXE+로컬 SQLite로 단절 없이 동작")
feat(11.7, 10.4, C_ACCENT, "🎯 SSOT — 단일 진실원",
     "법령(law-api)·정밀도 기준(Version11.xlsx)을 단일 소스로 주입 → 클라이언트·챗봇에 동일 복제 배포")
feat(22.5, 10.4, C_ANTI, "🛡 보안형 AI 개발",
     "승인 게이트 + 쓰기 범위 격리 + 컨텍스트 제한 + 감사 로그 → AI 코드 파괴 위험 원천 방지")

out = "KTL_통합시스템_설계.pptx"
prs.save(out)
print("저장:", out, "| 슬라이드", len(prs.slides._sldIdLst))
