from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── 색상 팔레트 ──────────────────────────────────────────
C_BG       = RGBColor(0x0F, 0x17, 0x2A)   # 배경 (진네이비)
C_DEV      = RGBColor(0x1E, 0x3A, 0x5F)   # 개발자 영역
C_ANTI     = RGBColor(0x2E, 0x4A, 0x7F)   # Antigravity
C_CLAUDE   = RGBColor(0xD9, 0x7B, 0x3A)   # Claude (주황)
C_MCP      = RGBColor(0x8B, 0x5C, 0xF6)   # MCP (보라)
C_VERCEL   = RGBColor(0x00, 0x00, 0x00)   # Vercel
C_MAC      = RGBColor(0x1A, 0x4A, 0x3A)   # Mac Studio (초록)
C_EXT      = RGBColor(0x2A, 0x2A, 0x4A)   # 외부 API
C_USER     = RGBColor(0x1A, 0x3A, 0x5A)   # 웹사용자
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_YELLOW   = RGBColor(0xFF, 0xD7, 0x00)
C_GRAY     = RGBColor(0x88, 0x99, 0xAA)
C_GREEN    = RGBColor(0x4A, 0xDE, 0x80)
C_ORANGE   = RGBColor(0xFF, 0x99, 0x44)
C_PURPLE   = RGBColor(0xC0, 0xA0, 0xFF)

def bg_rect(slide, clr):
    sh = slide.shapes.add_shape(1, Cm(0), Cm(0), prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = clr
    sh.line.fill.background()
    return sh

def box(slide, x, y, w, h, fill, line_clr=None, radius=0):
    sh = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_clr:
        sh.line.color.rgb = line_clr
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def label(slide, x, y, w, h, text, size=9, bold=False, clr=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = clr
    return tb

def connector(slide, x1, y1, x2, y2, clr=C_GRAY, width=1.5):
    from pptx.util import Cm as C
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    # 직선 커넥터
    cx = min(x1,x2); cy = min(y1,y2)
    cw = abs(x2-x1); ch = abs(y2-y1)
    if cw < 0.01: cw = 0.01
    if ch < 0.01: ch = 0.01
    conn = slide.shapes.add_connector(1, C(x1), C(y1), C(x2), C(y2))
    conn.line.color.rgb = clr
    conn.line.width = Pt(width)
    return conn

# ── 배경 ────────────────────────────────────────────────
bg_rect(slide, C_BG)

# ── 제목 ────────────────────────────────────────────────
label(slide, 0.5, 0.2, 32, 0.9,
      "KTL 정도검사 계산기 — 전체 시스템 구성도",
      size=16, bold=True, clr=C_WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
#  ROW 1: 개발자 영역
# ════════════════════════════════════════════════════════
# 개발자 영역 배경
box(slide, 0.4, 1.3, 32.9, 6.4, RGBColor(0x16, 0x24, 0x3E),
    line_clr=RGBColor(0x44, 0x66, 0x99))
label(slide, 0.7, 1.35, 6, 0.5, "개발 환경", size=7, clr=C_GRAY)

# Antigravity IDE
box(slide, 0.6, 1.9, 7.2, 5.3, C_ANTI, line_clr=RGBColor(0x66,0x99,0xFF))
label(slide, 0.8, 2.0, 6.8, 0.6, "Antigravity IDE  (안티)", size=10, bold=True, clr=C_YELLOW)
label(slide, 0.8, 2.7, 6.6, 4.0,
      "✦ 코드 편집 · 파일 탐색\n"
      "✦ 터미널 통합\n"
      "✦ Git 연동\n"
      "✦ Claude Code 연동\n"
      "✦ pyrefly (Python 분석)\n\n"
      "→  git push  →  Vercel 자동 배포",
      size=8.5, clr=C_WHITE)

# Claude Code
box(slide, 8.3, 1.9, 8.5, 2.4, C_CLAUDE, line_clr=RGBColor(0xFF,0xAA,0x55))
label(slide, 8.5, 2.0, 8.0, 0.6, "Claude Code  (AI 뇌)", size=10, bold=True, clr=C_WHITE)
label(slide, 8.5, 2.7, 8.0, 1.4,
      "✦ 자연어 ↔ 코드 변환\n"
      "✦ 도구 호출 판단\n"
      "✦ Vercel 배포 · 관리",
      size=8.5, clr=C_WHITE)

# MCP 서버
box(slide, 8.3, 4.5, 8.5, 2.8, C_MCP, line_clr=RGBColor(0xAA,0x77,0xFF))
label(slide, 8.5, 4.6, 8.0, 0.55, "ktl-calculator MCP  (12 Tools)", size=9.5, bold=True, clr=C_WHITE)
label(slide, 8.5, 5.2, 8.0, 1.9,
      "계산(1)  calculate_accuracy  오차율·합격 판정\n"
      "DB  (3)  list_items · get_fee · get_sheet\n"
      "법령(3)  search_laws · get_content · get_legal_basis\n"
      "지식(4)  search_knowledge · get_node · status",
      size=8, clr=C_WHITE)

# MCP 데이터소스 3개
# Excel DB
box(slide, 17.4, 1.9, 4.8, 2.3, RGBColor(0x1A,0x3A,0x2A), line_clr=RGBColor(0x44,0xAA,0x66))
label(slide, 17.6, 2.0, 4.4, 0.5, "Excel DB", size=9, bold=True, clr=C_GREEN)
label(slide, 17.6, 2.55, 4.4, 1.4,
      "Version11_(2026).xlsx\n수수료 · 기준 · 항목\n싱글턴 캐시",
      size=8, clr=C_WHITE)

# law.go.kr (MCP용)
box(slide, 17.4, 4.4, 4.8, 2.0, RGBColor(0x2A,0x1A,0x3A), line_clr=RGBColor(0x88,0x44,0xAA))
label(slide, 17.6, 4.5, 4.4, 0.5, "law.go.kr DRF", size=9, bold=True, clr=C_PURPLE)
label(slide, 17.6, 5.05, 4.4, 1.1,
      "XML 2.3 MB → 파싱 1.8 KB\n99.9% 토큰 절감 · 1일 캐시",
      size=8, clr=C_WHITE)

# Obsidian KB
box(slide, 22.8, 1.9, 5.2, 4.5, RGBColor(0x2A,0x2A,0x1A), line_clr=RGBColor(0xAA,0xAA,0x44))
label(slide, 23.0, 2.0, 4.8, 0.5, "Obsidian KB  (knowledge/)", size=9, bold=True, clr=C_YELLOW)
label(slide, 23.0, 2.6, 4.8, 3.5,
      "7개 카테고리\n"
      "  간이측정기 · 공통\n"
      "  법령 · 성능시험\n"
      "  정도검사 · 표준용액\n"
      "  형식승인\n\n"
      "35개+ .md 노드\n[[링크]] 그래프 탐색",
      size=8, clr=C_WHITE)

# 커넥터: Claude Code → MCP
connector(slide, 12.55, 4.3, 12.55, 4.5, clr=C_ORANGE, width=2)
# 커넥터: MCP → Excel
connector(slide, 16.8, 5.3, 17.4, 3.0, clr=C_GREEN, width=1.5)
# 커넥터: MCP → law.go.kr
connector(slide, 16.8, 5.7, 17.4, 5.4, clr=C_PURPLE, width=1.5)
# 커넥터: MCP → Obsidian
connector(slide, 16.8, 5.0, 22.8, 4.0, clr=C_YELLOW, width=1.5)

# ════════════════════════════════════════════════════════
#  ROW 2: Vercel
# ════════════════════════════════════════════════════════
box(slide, 0.4, 8.1, 21.5, 4.8, RGBColor(0x10,0x10,0x10),
    line_clr=RGBColor(0xFF,0xFF,0xFF))
label(slide, 0.7, 8.2, 8, 0.55, "Vercel  (icn1 · 서울)", size=10, bold=True, clr=C_WHITE)

# CDN
box(slide, 0.6, 8.85, 5.2, 3.6, RGBColor(0x1A,0x1A,0x2A), line_clr=RGBColor(0x55,0x55,0x88))
label(slide, 0.8, 8.95, 4.8, 0.5, "CDN  (정적)", size=9, bold=True, clr=RGBColor(0xAA,0xBB,0xFF))
label(slide, 0.8, 9.55, 4.8, 2.6,
      "dist/  (Vite 빌드)\nHTML · JS · CSS\n\n번들 포함:\n  Excel DB (.xlsx)\n  KB (.md × 35)",
      size=8, clr=C_WHITE)

# Serverless
box(slide, 6.3, 8.85, 15.2, 3.6, RGBColor(0x1A,0x1A,0x2A), line_clr=RGBColor(0x55,0x55,0x88))
label(slide, 6.5, 8.95, 14.8, 0.5, "Serverless Functions  (api/)", size=9, bold=True, clr=RGBColor(0xAA,0xBB,0xFF))
label(slide, 6.5, 9.55, 14.5, 2.6,
      "/api/auth          로그인 · 토큰 검증\n"
      "/api/lawChat       법령 AI 챗봇  (Obsidian + law.go.kr + Gemini)\n"
      "/api/gemini        범용 AI 프록시\n"
      "/api/lawSearch     법령명 검색  (law.go.kr XML)\n"
      "/api/legalBasis    항목별 정도검사 근거 · 기준\n"
      "/api/calcData      계산 결과 저장 · 불러오기  (→ Mac Studio)\n"
      "/api/db/status     Excel DB 상태  /  /api/admin  관리자",
      size=8, clr=C_WHITE)

# Vercel → 외부 연결
# Gemini
box(slide, 22.5, 8.1, 5.5, 2.3, RGBColor(0x1A,0x2A,0x3A), line_clr=RGBColor(0x44,0x88,0xCC))
label(slide, 22.7, 8.2, 5.1, 0.55, "Gemini 2.5 Flash", size=9, bold=True, clr=RGBColor(0x66,0xBB,0xFF))
label(slide, 22.7, 8.85, 5.1, 1.2, "Google AI API\nlawChat 답변 생성\n범용 AI 프록시", size=8, clr=C_WHITE)

# law.go.kr (Vercel용)
box(slide, 22.5, 10.6, 5.5, 2.3, RGBColor(0x1A,0x2A,0x3A), line_clr=RGBColor(0x44,0x88,0xCC))
label(slide, 22.7, 10.7, 5.1, 0.55, "law.go.kr  국가법령정보", size=9, bold=True, clr=RGBColor(0x66,0xBB,0xFF))
label(slide, 22.7, 11.35, 5.1, 1.2, "실시간 법령 조회\n조문 XML 수신", size=8, clr=C_WHITE)

connector(slide, 21.5, 9.2, 22.5, 9.2, clr=RGBColor(0x44,0x88,0xCC), width=1.5)
connector(slide, 21.5, 10.5, 22.5, 11.0, clr=RGBColor(0x44,0x88,0xCC), width=1.5)

# ════════════════════════════════════════════════════════
#  ROW 3: Mac Studio
# ════════════════════════════════════════════════════════
box(slide, 0.4, 13.3, 21.5, 4.8, C_MAC, line_clr=RGBColor(0x44,0xCC,0x77))
label(slide, 0.7, 13.4, 10, 0.55, "Mac Studio  (로컬 · pm2 상시 실행)", size=10, bold=True, clr=C_GREEN)

# parser-server
box(slide, 0.6, 14.05, 9.8, 3.6, RGBColor(0x12,0x32,0x22), line_clr=RGBColor(0x33,0x99,0x55))
label(slide, 0.8, 14.15, 9.4, 0.5, "parser-server  :3333", size=9, bold=True, clr=C_GREEN)
label(slide, 0.8, 14.75, 9.4, 2.6,
      "사진 OCR · 파싱\n작업 상태 관리\n대시보드",
      size=8.5, clr=C_WHITE)

# mac-studio-server
box(slide, 11.0, 14.05, 10.5, 3.6, RGBColor(0x12,0x32,0x22), line_clr=RGBColor(0x33,0x99,0x55))
label(slide, 11.2, 14.15, 10.1, 0.5, "mac-studio-server  :3456", size=9, bold=True, clr=C_GREEN)
label(slide, 11.2, 14.75, 10.1, 2.6,
      "calc_data  (SQLite DB)\n계산 결과 저장 · 불러오기\nlocations  현장 위치 조회",
      size=8.5, clr=C_WHITE)

# cloudflared 경고
box(slide, 22.5, 13.3, 5.5, 4.8, RGBColor(0x2A,0x1A,0x0A), line_clr=RGBColor(0xFF,0x99,0x00))
label(slide, 22.7, 13.4, 5.1, 0.55, "⚠️  cloudflared", size=9, bold=True, clr=C_ORANGE)
label(slide, 22.7, 14.05, 5.1, 3.8,
      "quick tunnel\n(재시작마다 URL 변경)\n\n현재:\n*.trycloudflare.com\n\n→ Named Tunnel\n  전환 권장",
      size=8, clr=C_WHITE)

# Vercel → Mac Studio 커넥터
connector(slide, 11.15, 12.9, 11.15, 13.3, clr=C_GREEN, width=2)

# ════════════════════════════════════════════════════════
#  웹 사용자 (우측 세로)
# ════════════════════════════════════════════════════════
box(slide, 28.6, 1.3, 4.8, 16.8, RGBColor(0x0A,0x22,0x3A), line_clr=RGBColor(0x44,0xAA,0xFF))
label(slide, 28.8, 1.5, 4.4, 0.55, "웹 사용자", size=10, bold=True, clr=RGBColor(0x66,0xCC,0xFF))
label(slide, 28.8, 2.2, 4.4, 0.5, "현장 엔지니어", size=8, clr=C_GRAY)
label(slide, 28.8, 3.0, 4.2, 14.5,
      "브라우저\n\ncalculator-\nsnowy-eight-\n87.vercel.app\n\n\n"
      "✦ 정도검사\n   계산\n\n"
      "✦ 법령 AI\n   챗봇\n\n"
      "✦ 계산 결과\n   저장·불러\n   오기\n\n"
      "✦ 성적서\n   Claydox\n   전송",
      size=8.5, clr=C_WHITE)

# 웹 사용자 → Vercel 커넥터
connector(slide, 28.6, 9.5, 21.9, 9.5, clr=RGBColor(0x44,0xAA,0xFF), width=2)

# ════════════════════════════════════════════════════════
#  범례
# ════════════════════════════════════════════════════════
box(slide, 0.4, 18.35, 28.0, 0.55, RGBColor(0x0A,0x0A,0x1A))
label(slide, 0.6, 18.4, 27.5, 0.45,
      "  🟠 Claude Code (AI)     🟣 MCP 12 Tools     ⚫ Vercel 서버리스     🟢 Mac Studio 로컬     🔵 웹 사용자 (HTTPS)     ⚠️ cloudflared 주의",
      size=7.5, clr=C_GRAY, align=PP_ALIGN.LEFT)

out = "/Users/kimjongjin/coding/calculator-main/KTL_시스템_구성도.pptx"
prs.save(out)
print("저장:", out)
