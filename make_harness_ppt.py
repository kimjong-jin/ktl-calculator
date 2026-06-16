# -*- coding: utf-8 -*-
"""KTL Workspace 멀티에이전트 하네스(Harness) 시스템 발표 PPT 생성."""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
SW, SH = 33.87, 19.05

# ── 팔레트 ───────────────────────────────────────────────
C_BG     = RGBColor(0x0B, 0x12, 0x22)
C_PANEL  = RGBColor(0x16, 0x21, 0x38)
C_ANTI   = RGBColor(0x6D, 0x5A, 0xE0)   # Antigravity 오케스트레이터 (보라)
C_CLAUDE = RGBColor(0xD9, 0x7B, 0x3A)   # Claude (주황)
C_CODEX  = RGBColor(0x1F, 0x9D, 0x8A)   # Codex (청록)
C_HUMAN  = RGBColor(0xC2, 0x96, 0x2C)   # 인간 승인 (골드)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_INK    = RGBColor(0x0B, 0x12, 0x22)
C_SUB    = RGBColor(0x9F, 0xB3, 0xC8)
C_LINE   = RGBColor(0x2C, 0x3C, 0x5A)
C_HEAD   = RGBColor(0x24, 0x33, 0x52)
C_ROW1   = RGBColor(0x12, 0x1C, 0x30)
C_ROW2   = RGBColor(0x16, 0x22, 0x3A)
C_ACCENT = RGBColor(0x4A, 0xDE, 0x80)

def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C_BG; r.line.fill.background()
    r.shadow.inherit = False
    return s

def box(s, x, y, w, h, fill, line=None, radius=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(s, x, y, w, h, text, size, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font='맑은 고딕', spacing=1.0):
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(0.1); tf.margin_top = tf.margin_bottom = Cm(0.05)
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

def label_box(s, x, y, w, h, fill, title, sub, tcolor=C_WHITE):
    box(s, x, y, w, h, fill, radius=True)
    txt(s, x, y+0.35, w, 1.0, title, 13, tcolor, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    txt(s, x, y+h-1.15, w, 1.0, sub, 9.5, tcolor, False, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ════════════════════ 슬라이드 1: 타이틀 ════════════════════
s = slide()
box(s, 0, 0, 1.0, SH, C_ANTI)
txt(s, 2.4, 4.3, 29, 1.2, "KTL WORKSPACE", 20, C_SUB, True)
txt(s, 2.4, 5.6, 29, 3.6,
    "멀티에이전트 하네스(Harness) 시스템", 40, C_WHITE, True, spacing=1.0)
txt(s, 2.4, 9.2, 29, 1.5,
    "차세대 보안형 AI 개발 프레임워크", 22, C_CLAUDE, True)
# 칩
chips = [("오케스트레이터 · 워커풀", C_ANTI), ("인간 승인 게이트", C_HUMAN),
         ("쓰기 범위 격리", C_CODEX), ("Append-Only 감사로그", C_ACCENT)]
cx = 2.4
for t, c in chips:
    w = 0.55 + len(t) * 0.42
    b = box(s, cx, 12.0, w, 1.1, C_PANEL, line=c, radius=True)
    txt(s, cx, 12.0, w, 1.1, t, 11, C_WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    cx += w + 0.5
txt(s, 2.4, 17.4, 29, 1.0,
    "AI에게 맹목적으로 맡기지 않는다 — 인간의 통제 하에 협업하는 AI 개발 인프라",
    13, C_SUB, False)

# ════════════════════ 슬라이드 2: 역할 아키텍처 ════════════════════
s = slide()
txt(s, 1.2, 0.8, 31, 1.2, "역할 기반 아키텍처 (Role-based Architecture)", 24, C_WHITE, True)
txt(s, 1.2, 2.1, 31, 0.9, "기획은 오케스트레이터가, 변경·코딩은 전문 워커가 — 인간 승인 게이트 통과 후 실행", 13, C_SUB)

# 인간
label_box(s, 13.4, 3.4, 7.0, 1.7, C_HUMAN,
          "👤  인간 개발자 (승인 게이트)", "명시적 승인 없이는 어떤 외부 AI도 자동 실행 불가", C_INK)
# 오케스트레이터
label_box(s, 12.4, 6.2, 9.0, 2.4, C_ANTI,
          "🧭  오케스트레이터  ·  Antigravity (Gemini 3.1 Pro High)",
          "전체 컨텍스트 기획 · 마일스톤 제어 · 긴 문서 분석으로 계획 수립")
# 디스패처
box(s, 14.9, 9.1, 4.1, 1.1, C_PANEL, line=C_LINE, radius=True)
txt(s, 14.9, 9.1, 4.1, 1.1, "call_worker.sh  (API 토큰 최소 노출)", 10, C_SUB, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
# 워커
label_box(s, 7.0, 11.0, 9.2, 2.3, C_CLAUDE,
          "🟧  워커 · claude-main", "실제 코드 변경 · 작성 (추론·구현 강점)", C_INK)
label_box(s, 17.7, 11.0, 9.2, 2.3, C_CODEX,
          "🟦  워커 · codex-main", "실제 코드 변경 · 작성 (코딩 강점)", C_WHITE)
# 워커풀 라벨
txt(s, 7.0, 13.5, 19.9, 0.8, "워커 풀 (Worker Pool) — 본인 작업영역(tasks/) 내에서만 작성", 11, C_SUB, True, PP_ALIGN.CENTER)

# 화살표 (간단 라인)
def vline(x, y1, y2, c=C_LINE):
    ln = s.shapes.add_connector(2, Cm(x), Cm(y1), Cm(x), Cm(y2))
    ln.line.color.rgb = c; ln.line.width = Pt(2.5)
vline(16.9, 5.1, 6.2, C_HUMAN)
vline(16.9, 8.6, 9.1, C_ANTI)
vline(11.6, 10.2, 11.0, C_CLAUDE); vline(22.3, 10.2, 11.0, C_CODEX)
hl = s.shapes.add_connector(2, Cm(11.6), Cm(10.2), Cm(22.3), Cm(10.2)); hl.line.color.rgb = C_LINE; hl.line.width = Pt(2.5)

txt(s, 1.2, 16.7, 31, 1.4,
    "기대효과  ·  에이전트 간 역할 충돌 방지 + 각 모델의 강점(추론/코딩) 극대화 → 비용·결과물 완성도 최적화",
    13, C_ACCENT, True)

# ════════════════════ 슬라이드 3: 운영 요약표 ════════════════════
s = slide()
txt(s, 1.0, 0.7, 31, 1.1, "하네스 시스템 운영 요약", 24, C_WHITE, True)
txt(s, 1.0, 1.9, 31, 0.8, "5대 통제 정책 · 구현 검증 방식 · 비즈니스 기대효과", 13, C_SUB)

rows_data = [
    ("영역", "통제 정책", "구현 · 검증 방식", "비즈니스 기대효과"),
    ("에이전트\n역할 분담",
     "오케스트레이터=기획·마일스톤 / 워커풀(Claude·Codex)=코드 작성에 위임",
     "Gemini 3.1 Pro 문서분석으로 계획 · call_worker.sh 디스패처로 토큰 최소 노출",
     "역할 충돌 방지 · 모델별 강점(추론/코딩) 극대화 → 비용·완성도 최적화"),
    ("엄격한\n승인 게이트",
     "인간의 명시적 승인 없이는 어떤 외부 AI 엔진도 자동 실행 불가",
     "task.md의 workers_approved 기록 · 유료 API 호출 전 비용 알림 + 터미널 승인 대기",
     "오버빌링(과금 폭탄) 방지 · 에이전트 오작동/폭주 조기 안전 통제"),
    ("코드 쓰기\n범위 통제",
     "워커는 본인 작업영역(tasks/) 외 프로젝트 소스를 임의 변경하지 못하도록 격리",
     "외부 쓰기 4조건: target_repo 명시 · write_scope 제한 · task.md 기입 · [APPROVAL] 로그",
     "코드 유실 위험 차단 · 개발 브랜치 안정성(Stability) 보장"),
    ("컨텍스트\n크기 제어",
     "히스토리 비대화로 모델이 환각(Hallucination) 일으키는 현상 방지",
     "wc -m / wc -w 상시 측정 · context.md ≤ 1500자 · brief.md ≤ 1200자 정량 제한",
     "프롬프트 최소화 → 반응 속도↑ · 오작동·오답률 최소화"),
    ("이력 기록\n· 사후 검증",
     "모든 중대 결정·호출·검증을 수정/삭제 불가 로그로 영구 기록",
     "log.md Append-Only · 결과 획득 후 Verification Checklist → [VERIFICATION] 태그",
     "투명한 추적성(Traceability) 확보 · 감사(Audit) 가능 인프라"),
]
TX, TY = 1.0, 3.1
TW = 31.87
col_w = [4.6, 8.6, 9.8, 8.87]
header_h = 1.1
row_h = (SH - TY - 0.9 - header_h) / 5
y = TY
for ri, row in enumerate(rows_data):
    h = header_h if ri == 0 else row_h
    x = TX
    for ci, cell in enumerate(row):
        w = col_w[ci]
        if ri == 0:
            box(s, x, y, w, h, C_HEAD, line=C_BG)
            txt(s, x+0.25, y, w-0.5, h, cell, 13, C_WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        else:
            fill = C_ROW1 if ri % 2 else C_ROW2
            box(s, x, y, w, h, fill, line=C_BG)
            if ci == 0:
                # 영역 컬럼: 좌측 컬러바 + 강조
                bar_c = [C_ANTI, C_HUMAN, C_CODEX, C_CLAUDE, C_ACCENT][ri-1]
                box(s, x, y, 0.18, h, bar_c)
                txt(s, x+0.35, y, w-0.55, h, cell, 12.5, C_WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            else:
                col = C_ACCENT if ci == 3 else C_WHITE
                bold = ci == 3
                txt(s, x+0.28, y, w-0.5, h, cell, 10.5, col if ci==3 else C_SUB if ci==2 else C_WHITE,
                    bold, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, spacing=1.02)
        x += w
    y += h

# ════════════════════ 슬라이드 4: 핵심 멘트 ════════════════════
s = slide()
box(s, 0, 0, SW, SH, C_BG)
box(s, 2.0, 2.4, 0.6, 14.2, C_ANTI)
txt(s, 3.4, 2.6, 28, 1.2, "🎤  발표 핵심 멘트", 20, C_CLAUDE, True)
quote = ("“KTL Workspace에 구축된 멀티에이전트 하네스(Harness) 시스템은\n"
         "AI에게 작업을 맹목적으로 맡기는 것이 아니라,\n"
         "오케스트레이터(Antigravity)와 역할별 전문 워커(Claude·Codex)를\n"
         "파일 기반 메모리 시스템 위에 격리하여 협업하게 만듭니다.\n\n"
         "인간 개발자의 승인 게이트와 엄격한 쓰기 범위 제한 하에 구동되므로,\n"
         "AI의 코드 파괴 위험을 원천 방지하며 고도의 품질을 담보하는\n"
         "차세대 보안형 AI 개발 프레임워크입니다.”")
txt(s, 3.4, 4.4, 28.5, 10.5, quote, 20, C_WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, spacing=1.25)
# 하단 키워드
keys = ["통제(Control)", "격리(Isolation)", "추적성(Traceability)", "안전(Safety)"]
cx = 3.4
for k in keys:
    w = 0.6 + len(k) * 0.42
    box(s, cx, 16.4, w, 1.1, C_PANEL, line=C_ACCENT, radius=True)
    txt(s, cx, 16.4, w, 1.1, k, 12, C_ACCENT, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    cx += w + 0.5

out = "KTL_하네스_시스템.pptx"
prs.save(out)
print("저장:", out, "| 슬라이드", len(prs.slides._sldIdLst))
